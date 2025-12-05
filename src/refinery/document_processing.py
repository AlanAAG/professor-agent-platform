# src/refinery/document_processing.py

import logging
import os
from typing import Optional
from urllib.parse import quote

# Lazy imports for local office file processing to avoid hard dependencies if not used
try:
    import docx
except ImportError:
    docx = None

try:
    import pptx
except ImportError:
    pptx = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

from selenium import webdriver
from selenium.webdriver.common.by import By
from selenium.webdriver.support.ui import WebDriverWait
from selenium.webdriver.support import expected_conditions as EC
from selenium.common.exceptions import TimeoutException, WebDriverException


def _navigates_to_login(current_url: str) -> bool:
    lower = (current_url or "").lower()
    return any(k in lower for k in ["accounts.google.com", "/signin", "/login", "service=wise", "auth" ])


def extract_google_docs_content(driver: webdriver.Chrome, url: str) -> Optional[str]:
    """Navigate to Google Docs and extract visible text from editor."""
    try:
        driver.get(url)
        WebDriverWait(driver, 30).until(EC.presence_of_element_located((By.CSS_SELECTOR, ".kix-appview-editor")))
        if _navigates_to_login(driver.current_url):
            logging.warning("Redirected to login while accessing Google Docs")
            return None
        pages = driver.find_elements(By.CSS_SELECTOR, ".kix-page")
        texts = [p.text.strip() for p in pages if getattr(p, "text", "").strip()]
        content = "\n\n".join(texts) if texts else "(Empty document)"
        logging.info(f"Extracted Google Docs content ({len(content)} chars)")
        return content
    except (TimeoutException, WebDriverException) as e:
        logging.warning(f"Failed to extract Google Docs content: {e}")
        return None


def extract_google_sheets_content(driver: webdriver.Chrome, url: str) -> Optional[str]:
    """Extract visible cell data from Google Sheets as tab-separated text."""
    try:
        driver.get(url)
        WebDriverWait(driver, 30).until(EC.presence_of_element_located((By.CSS_SELECTOR, ".docs-sheet-tab-panel")))
        if _navigates_to_login(driver.current_url):
            logging.warning("Redirected to login while accessing Google Sheets")
            return None
        # Visible table grid
        rows = driver.find_elements(By.CSS_SELECTOR, "table tr")
        lines = []
        max_rows = 100  # limit very large sheets
        truncated = False
        for idx, row in enumerate(rows):
            if idx >= max_rows:
                truncated = True
                break
            cells = row.find_elements(By.CSS_SELECTOR, "td")
            line = "\t".join([c.text.replace("\n", " ").strip() for c in cells])
            if line.strip():
                lines.append(line)
        content = "\n".join(lines) if lines else "(Empty document)"
        if truncated:
            content += "\n\n(Truncated: showing first 100 rows)"
        logging.info(f"Extracted Google Sheets content ({len(content)} chars)")
        return content
    except (TimeoutException, WebDriverException) as e:
        logging.warning(f"Failed to extract Google Sheets content: {e}")
        return None


def extract_google_slides_content(driver: webdriver.Chrome, url: str) -> Optional[str]:
    """Extract slide text content from Google Slides."""
    try:
        driver.get(url)
        WebDriverWait(driver, 30).until(EC.presence_of_element_located((By.CSS_SELECTOR, ".punch-viewer-page")))
        if _navigates_to_login(driver.current_url):
            logging.warning("Redirected to login while accessing Google Slides")
            return None
        # Gather text elements from slide content containers
        slides = driver.find_elements(By.CSS_SELECTOR, ".punch-viewer-page")
        max_slides = 20
        parts = []
        for i, slide in enumerate(slides[:max_slides]):
            # Common containers for text content
            texts = []
            for sel in [
                ".punch-viewer-content",
                ".punch-viewewr-slide-text",  # typo variant per spec
                "div[role='region']",
            ]:
                try:
                    nodes = slide.find_elements(By.CSS_SELECTOR, sel)
                    for n in nodes:
                        t = n.text.strip()
                        if t:
                            texts.append(t)
                except Exception:
                    continue
            joined = "\n".join(texts) if texts else "(No text on slide)"
            parts.append(f"Slide {i+1}:\n{joined}")
        if len(slides) > max_slides:
            parts.append("(Truncated: showing first 20 slides)")
        content = "\n\n".join(parts) if parts else "(Empty document)"
        logging.info(f"Extracted Google Slides content ({len(content)} chars)")
        return content
    except (TimeoutException, WebDriverException) as e:
        logging.warning(f"Failed to extract Google Slides content: {e}")
        return None


def extract_office_online_content(driver: webdriver.Chrome, original_url: str) -> Optional[str]:
    """Wrap any office file URL in Office Online viewer and extract content."""
    try:
        viewer_url = f"https://view.officeapps.live.com/op/view.aspx?src={quote(original_url)}"
        driver.get(viewer_url)
        WebDriverWait(driver, 30).until(EC.presence_of_element_located((By.CSS_SELECTOR, "#WACViewPanel")))
        if _navigates_to_login(driver.current_url):
            logging.warning("Redirected to login while accessing Office Online viewer")
            return None
        # Common content containers
        nodes = driver.find_elements(By.CSS_SELECTOR, ".WACViewPanel_Content, div[role='document']")
        texts = [n.text.strip() for n in nodes if getattr(n, "text", "").strip()]
        content = "\n\n".join(texts) if texts else "(Empty document)"
        logging.info(f"Extracted Office Online content ({len(content)} chars)")
        return content
    except (TimeoutException, WebDriverException) as e:
        logging.warning(f"Failed to extract Office Online content: {e}")
        return None


def extract_drive_content(driver: webdriver.Chrome, url: str, doc_type: str) -> Optional[str]:
    """
    Universal extractor that routes to appropriate viewer based on URL/type.

    Args:
        driver: Selenium WebDriver instance
        url: Document URL (Google Docs/Sheets/Slides or direct office file)
        doc_type: One of: "google_docs", "google_sheets", "google_slides", "office_document"

    Returns:
        Extracted text content or None if extraction fails
    """
    try:
        doc_type = (doc_type or "").lower()
        if doc_type == "google_docs" or "docs.google.com/document" in url:
            return extract_google_docs_content(driver, url)
        if doc_type == "google_sheets" or "docs.google.com/spreadsheets" in url:
            return extract_google_sheets_content(driver, url)
        if doc_type == "google_slides" or "docs.google.com/presentation" in url:
            return extract_google_slides_content(driver, url)
        # Default: wrap via Office Online viewer for office files
        return extract_office_online_content(driver, url)
    except Exception as e:
        logging.warning(f"Document extraction routing failed: {e}")
        return None


def process_local_office_file(file_path: str) -> str:
    """
    Extract text content from a local Office file (.docx, .pptx, .xlsx).

    Args:
        file_path: Path to the local file.

    Returns:
        Extracted text content or an empty string if processing fails or format is unsupported.
    """
    if not os.path.exists(file_path):
        logging.error(f"File not found: {file_path}")
        return ""

    file_ext = os.path.splitext(file_path)[1].lower()
    text_content = ""

    try:
        if file_ext == ".docx":
            if docx is None:
                logging.error("python-docx is not installed.")
                return ""
            try:
                doc = docx.Document(file_path)
                # Extract text from paragraphs
                full_text = []
                for para in doc.paragraphs:
                    full_text.append(para.text)
                text_content = "\n".join(full_text)
            except Exception as e:
                logging.error(f"Failed to process DOCX {file_path}: {e}")
                return ""

        elif file_ext == ".pptx":
            if pptx is None:
                logging.error("python-pptx is not installed.")
                return ""
            try:
                prs = pptx.Presentation(file_path)
                text_runs = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text_runs.append(shape.text)
                text_content = "\n".join(text_runs)
            except Exception as e:
                logging.error(f"Failed to process PPTX {file_path}: {e}")
                return ""

        elif file_ext == ".xlsx":
            if openpyxl is None:
                logging.error("openpyxl is not installed.")
                return ""
            try:
                wb = openpyxl.load_workbook(file_path, data_only=True)
                # Active sheet only? Prompt says "all cells in the active sheet".
                sheet = wb.active
                if sheet:
                    rows = []
                    for row in sheet.iter_rows(values_only=True):
                        # Filter None values and join
                        row_text = "\t".join([str(cell) for cell in row if cell is not None])
                        if row_text.strip():
                            rows.append(row_text)
                    text_content = "\n".join(rows)
            except Exception as e:
                logging.error(f"Failed to process XLSX {file_path}: {e}")
                return ""

        else:
             logging.warning(f"Unsupported office file extension: {file_ext}")
             return ""

        return text_content

    except Exception as e:
         logging.error(f"Unexpected error processing local office file {file_path}: {e}")
         return ""
